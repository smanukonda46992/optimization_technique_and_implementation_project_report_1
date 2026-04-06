from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
PPT_PATH = ROOT / "Part2_Elementary_Data_Structures_Presentation.pptx"
AUTHOR_NAME = "Sai Mohan Manukonda"
UC_STUDENT_ID = "005046992"


def ensure_title_slide_identity(prs):
    if not prs.slides:
        return

    slide = prs.slides[0]
    if len(slide.placeholders) < 2:
        return

    subtitle = slide.placeholders[1].text_frame
    subtitle.clear()
    subtitle.text = (
        "Assignment 6 Presentation\n"
        "Arrays, Stacks, Queues, Linked Lists, and Rooted Trees\n"
        f"{AUTHOR_NAME} | UC Student ID: {UC_STUDENT_ID} | Algorithms and Data Structures | April 2026"
    )


def main():
    if not PPT_PATH.exists():
        raise FileNotFoundError(f"PPT not found: {PPT_PATH}")

    prs = Presentation(str(PPT_PATH))
    prs.core_properties.author = AUTHOR_NAME
    prs.core_properties.last_modified_by = AUTHOR_NAME
    prs.core_properties.title = "Part 2 Elementary Data Structures Presentation"

    ensure_title_slide_identity(prs)
    prs.save(str(PPT_PATH))
    print(f"Updated author metadata and title identity in: {PPT_PATH}")


if __name__ == "__main__":
    main()
